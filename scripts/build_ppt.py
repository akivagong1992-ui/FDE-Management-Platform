"""生成 FDE 系统汇报 PowerPoint (模板风 + 精简文字)。

设计：
- 16:9 (13.33" × 7.5")
- 主色 #0066B3 (电信蓝) + 辅色 #00BCD4 (青) + 强调 #FF4081 (粉)
- 标题左上 + 蓝色装饰条 + 大字号 bullet
- 封面 / 章节扉页 / 内容页 3 种版式
- 文字精简：每页最多 5-6 个 bullet
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ─── 配色 ───────────────────────────────────────────────────
BLUE = RGBColor(0x00, 0x66, 0xB3)        # 电信主蓝
CYAN = RGBColor(0x00, 0xBC, 0xD4)        # 辅色青
PINK = RGBColor(0xFF, 0x40, 0x81)        # 强调粉
GOLD = RGBColor(0xFF, 0xB7, 0x40)        # 强调金
DARK = RGBColor(0x30, 0x31, 0x33)        # 正文深灰
GRAY = RGBColor(0x90, 0x93, 0x99)        # 副文字
LIGHT_BG = RGBColor(0xEC, 0xF5, 0xFF)    # 浅蓝背景
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

# ─── 尺寸 (16:9) ────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

BLANK_LAYOUT = prs.slide_layouts[6]  # 完全空白


# ─── 工具函数 ───────────────────────────────────────────────

def add_rect(slide, x, y, w, h, fill_color, line_color=None):
    """加一个矩形装饰条 / 背景块"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if line_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line_color
    shape.shadow.inherit = False
    return shape


def add_text(slide, x, y, w, h, text, *, font_size=18, bold=False,
             color=DARK, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             font_name="PingFang SC"):
    """加一段文字"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name
    return tb


def add_bullets(slide, x, y, w, h, items, *, font_size=18,
                bullet_color=BLUE, text_color=DARK, line_spacing=1.4,
                font_name="PingFang SC"):
    """加一组 bullet (用 ● 符号 + 缩进)"""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet 圆点
        run_b = p.add_run()
        run_b.text = "● "
        run_b.font.size = Pt(font_size)
        run_b.font.bold = True
        run_b.font.color.rgb = bullet_color
        run_b.font.name = font_name
        # 正文
        run_t = p.add_run()
        run_t.text = item
        run_t.font.size = Pt(font_size)
        run_t.font.color.rgb = text_color
        run_t.font.name = font_name
    return tb


def add_page_header(slide, page_num, total, title=None):
    """加页眉装饰条 + 标题"""
    # 左上角蓝色竖条装饰
    add_rect(slide, Inches(0.5), Inches(0.5), Inches(0.15), Inches(0.6), BLUE)
    if title:
        add_text(
            slide, Inches(0.75), Inches(0.45), Inches(10), Inches(0.7),
            title, font_size=28, bold=True, color=BLUE,
        )
    # 页眉右侧小字
    add_text(
        slide, Inches(10.5), Inches(0.5), Inches(2.5), Inches(0.4),
        f"FDE 团队管理系统 · v0.6 beta",
        font_size=10, color=GRAY, align=PP_ALIGN.RIGHT,
    )
    # 标题下分隔线
    add_rect(slide, Inches(0.5), Inches(1.2), Inches(12.33), Emu(20000), BLUE)


def add_page_footer(slide, page_num, total):
    """加页脚 + 页码"""
    add_text(
        slide, Inches(0.5), Inches(7.05), Inches(8), Inches(0.3),
        "中国电信国际香港分公司 · FDE 团队",
        font_size=10, color=GRAY,
    )
    add_text(
        slide, Inches(10.5), Inches(7.05), Inches(2.5), Inches(0.3),
        f"{page_num} / {total}",
        font_size=10, color=GRAY, align=PP_ALIGN.RIGHT,
    )


# ─── 版式构建器 ────────────────────────────────────────────

def make_cover():
    """封面 - 渐变蓝底 + 标题居中"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    # 渐变背景（用大矩形 + 渐变填充）
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    # 渐变填充用 lxml
    spPr = bg.fill._xPr.find(qn_safe('a:gradFill')) if False else None
    fill = bg.fill
    # 简单单色蓝（python-pptx 渐变不直接支持，用单色 + 装饰）
    fill.solid()
    fill.fore_color.rgb = BLUE
    bg.shadow.inherit = False

    # 顶部装饰条
    add_rect(s, 0, 0, SLIDE_W, Inches(0.15), CYAN)
    add_rect(s, 0, SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), CYAN)

    # 右下角装饰小方块
    add_rect(s, Inches(11), Inches(6.2), Inches(0.5), Inches(0.5), CYAN)
    add_rect(s, Inches(11.6), Inches(6.2), Inches(0.5), Inches(0.5), PINK)
    add_rect(s, Inches(12.2), Inches(6.2), Inches(0.5), Inches(0.5), GOLD)

    # 主标题
    add_text(s, Inches(1), Inches(2.5), Inches(11.33), Inches(1.2),
             "FDE 团队管理系统", font_size=60, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # 副标题
    add_text(s, Inches(1), Inches(3.8), Inches(11.33), Inches(0.8),
             "从外包黑盒到内化运营 · 数字化落地汇报",
             font_size=24, color=CYAN, align=PP_ALIGN.CENTER)
    # 底部信息
    add_text(s, Inches(1), Inches(5.3), Inches(11.33), Inches(0.5),
             "中国电信国际香港分公司 · FDE 团队",
             font_size=16, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(5.8), Inches(11.33), Inches(0.4),
             "v0.6-beta  ·  2026-05-24",
             font_size=12, color=LIGHT_BG, align=PP_ALIGN.CENTER)


def make_section(title, subtitle, page, total):
    """章节扉页 - 大编号 + 章节名"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    # 浅蓝背景
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, LIGHT_BG)
    # 左侧装饰大块
    add_rect(s, 0, 0, Inches(4.5), SLIDE_H, BLUE)
    # 大数字
    add_text(s, Inches(0.5), Inches(2.5), Inches(3.5), Inches(2),
             subtitle, font_size=120, bold=True, color=WHITE,
             align=PP_ALIGN.CENTER)
    # 标题
    add_text(s, Inches(5), Inches(3.0), Inches(7.8), Inches(1),
             title, font_size=44, bold=True, color=BLUE)
    # 小分隔线
    add_rect(s, Inches(5), Inches(4.0), Inches(1.2), Emu(40000), PINK)
    add_page_footer(s, page, total)


def make_content(page, total, title, items, *, subtitle=None,
                 notes=None, brag_box=None, item_size=18):
    """内容页 - 标题 + bullet 列表 + 可选讲话备注"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_page_header(s, page, total, title)
    y = Inches(1.5)
    if subtitle:
        add_text(s, Inches(0.75), y, Inches(12), Inches(0.5),
                 subtitle, font_size=16, color=GRAY)
        y = Inches(2.05)
    add_bullets(s, Inches(0.75), y, Inches(12), Inches(4.5),
                items, font_size=item_size)
    if brag_box:
        # 高亮关键 message 框
        box_y = Inches(5.3)
        add_rect(s, Inches(0.75), box_y, Inches(11.83), Inches(0.9),
                 LIGHT_BG, line_color=BLUE)
        add_text(s, Inches(1), box_y + Inches(0.15), Inches(11.5), Inches(0.6),
                 "⭐ " + brag_box, font_size=14, bold=True, color=BLUE,
                 anchor=MSO_ANCHOR.MIDDLE)
    if notes:
        # 讲话备注作为 speaker notes（演讲者视图才看得到，领导看不到）
        s.notes_slide.notes_text_frame.text = notes
    add_page_footer(s, page, total)


def make_table_page(page, total, title, headers, rows, *, notes=None,
                    col_widths=None, first_col_bold=True):
    """内容页 - 标题 + 表格"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    add_page_header(s, page, total, title)

    rows_n = len(rows) + 1
    cols_n = len(headers)
    tbl_left = Inches(0.75)
    tbl_top = Inches(1.6)
    tbl_w = Inches(11.83)
    tbl_h = Inches(0.5) + Inches(0.55) * len(rows)
    if tbl_h > Inches(5):
        tbl_h = Inches(5)

    tbl_shape = s.shapes.add_table(rows_n, cols_n, tbl_left, tbl_top, tbl_w, tbl_h)
    tbl = tbl_shape.table

    # 设置列宽
    if col_widths:
        total_emu = tbl_w
        sum_w = sum(col_widths)
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = int(total_emu * cw / sum_w)

    # 表头
    for j, h in enumerate(headers):
        cell = tbl.cell(0, j)
        cell.fill.solid()
        cell.fill.fore_color.rgb = BLUE
        cell.text = ""
        tf = cell.text_frame
        tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
        tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = h
        r.font.size = Pt(14); r.font.bold = True
        r.font.color.rgb = WHITE
        r.font.name = "PingFang SC"

    # 数据行
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = tbl.cell(i + 1, j)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if i % 2 == 0 else LIGHT_BG
            cell.text = ""
            tf = cell.text_frame
            tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
            tf.margin_top = Inches(0.05); tf.margin_bottom = Inches(0.05)
            p = tf.paragraphs[0]
            r = p.add_run()
            r.text = str(val)
            r.font.size = Pt(13)
            r.font.color.rgb = DARK
            r.font.name = "PingFang SC"
            if j == 0 and first_col_bold:
                r.font.bold = True
                r.font.color.rgb = BLUE

    if notes:
        s.notes_slide.notes_text_frame.text = notes
    add_page_footer(s, page, total)


def make_closing(page, total):
    """结语页 - 大字 + 决策问题"""
    s = prs.slides.add_slide(BLANK_LAYOUT)
    # 渐变背景
    add_rect(s, 0, 0, SLIDE_W, SLIDE_H, BLUE)
    add_rect(s, 0, 0, SLIDE_W, Inches(0.15), CYAN)
    add_rect(s, 0, SLIDE_H - Inches(0.15), SLIDE_W, Inches(0.15), CYAN)

    # 标语
    add_text(s, Inches(1), Inches(1.2), Inches(11.33), Inches(0.8),
             "系统化 · 透明化 · 合规化",
             font_size=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(1), Inches(2.1), Inches(11.33), Inches(0.5),
             "FDE 团队从「外包成本中心」向「公司战略资产」演进",
             font_size=18, color=CYAN, align=PP_ALIGN.CENTER)

    # 请领导决策标题
    add_text(s, Inches(2), Inches(3.3), Inches(9.33), Inches(0.6),
             "请领导决策",
             font_size=24, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # 2 个问题
    add_text(s, Inches(2), Inches(4.2), Inches(9.33), Inches(0.7),
             "1. 是否同意按 5 大战略方向推动业务扩展？",
             font_size=20, color=WHITE, align=PP_ALIGN.CENTER)
    add_text(s, Inches(2), Inches(5.0), Inches(9.33), Inches(0.7),
             "2. 系统正式上线天翼云 ECS 的预算 / 时间窗口？",
             font_size=20, color=WHITE, align=PP_ALIGN.CENTER)

    # 谢谢
    add_text(s, Inches(1), Inches(6.3), Inches(11.33), Inches(0.5),
             "—— 谢谢 ——",
             font_size=20, color=CYAN, align=PP_ALIGN.CENTER)


# 辅助 qn (避免 import 错误)
def qn_safe(tag): return tag


# ─── 内容 ───────────────────────────────────────────────────

TOTAL = 19

# 第 1 页 封面
make_cover()

# 第 2 页 议程
make_content(
    page=2, total=TOTAL,
    title="议程",
    subtitle="6 个部分 · 约 25 分钟",
    items=[
        "一、痛点 — 为什么需要这套系统",
        "二、解决方案 — 5 大支柱回应痛点",
        "三、系统化落地 — 8 维度业务闭环",
        "四、价值放大 — 系统能输出什么",
        "五、多角色协作 — 6 角色 × 11 模块",
        "六、未来演进 — 5 大战略方向",
    ],
    item_size=22,
    notes="今天汇报分 6 部分，前 2 块讲为什么 + 怎么解；中间 3 块讲怎么落地、谁在用、能产生什么价值；最后讲下一步方向。",
)

# 第 3 页 痛点黑盒
make_table_page(
    page=3, total=TOTAL,
    title="一、痛点：服务商模式的 4 个黑盒",
    headers=["黑盒", "现象", "后果"],
    rows=[
        ["💸 成本黑盒", "Vendor 报价多少给多少，无对比依据", "持续被报高价，议价无据"],
        ["👥 人力黑盒", "工程师能力 / 排班 / 满负荷度全在 Vendor 手上", "关键人员流失自己不知道"],
        ["📊 数据黑盒", "项目交付完资料散落，知识无沉淀", "同类项目重复踩坑"],
        ["🎯 利润瓜分", "项目利润大头被服务商截留，公司只拿薄利", "公司付了大钱，团队没成长，利润给了别人"],
    ],
    col_widths=[2.5, 5, 4.5],
    notes="这 4 个黑盒是过去一直存在但说不清的痛。系统建设目标就是打开这 4 个黑盒。利润瓜分是最痛的一条 — 公司付了钱，但赚的最多的不是公司。",
)

# 第 4 页 量化痛点对比
make_table_page(
    page=4, total=TOTAL,
    title="一、痛点：能看见的差距",
    headers=["数据点", "服务商模式", "FDE 模式"],
    rows=[
        ["客户付了多少", "✅ 知道（合同）", "✅ 知道"],
        ["Vendor 报价多少", "✅ 知道", "✅ 知道"],
        ["团队真实做了什么活", "❌ Vendor 报告", "✅ 工时数据"],
        ["工程师能力等级", "❌ 合同写 '高级工程师'", "✅ 厂商认证 L1/L2/L3"],
        ["客户满意度", "⚠️ 口头反馈", "✅ 1-5 星复盘"],
    ],
    col_widths=[3, 4, 5],
    notes="以前对一个项目的认知停留在合同金额。系统上线后从立项到复盘每个数据点都有据可查。",
)

# 第 5 页 5 大支柱总览
make_content(
    page=5, total=TOTAL,
    title="二、解决方案：5 大支柱",
    subtitle="5 个支柱 · 系统性回应 4 个黑盒",
    items=[
        "① 外部工程师高效管理 — 建档 + 排班 + 利用率",
        "② 技术能力管理沉淀 — 认证 + 培训 + IDP + 知识库",
        "③ 利润清晰可控 — 三口径并存，按角色分发",
        "④ 团队成本精细管理 — 8 类支出 + 审批闭环",
        "⑤ 驾驶舱价值展示 — 5 Tab 大屏，对外讲故事",
    ],
    item_size=22,
    brag_box="每个支柱回应至少 1 个黑盒；5 个支柱组合 = 系统的完整价值",
    notes="解决方案拆 5 个支柱。下面 5 页每页一个。",
)

# 第 6 页 支柱 ① 外部工程师高效管理
make_content(
    page=6, total=TOTAL,
    title="支柱 ① · 外部工程师高效管理",
    subtitle="合规建档 + 系统化排班 + 资源最大化",
    items=[
        "工程师独立建档：脱离 Vendor 视角，作为团队资源管理",
        "2 种签约形态：Vendor 直签 / Vendor 通过劳务公司",
        "证件号脱敏 + 解密审计：默认掩码，明文留痕",
        "派单 Assignment：PM 直接派工，接拒留对话",
        "工时按香港规则加权：工作日 1.0× / 加班 1.5×",
        "满负荷率实时计算 → 驾驶舱「Top 满负荷工程师」榜",
    ],
    item_size=18,
    brag_box="主导权回到团队：从 'Vendor 派人来' 变成 '我们派 Vendor 派的人'",
    notes="外包路线不放弃 — 这是合规和灵活性选择。但通过系统每个工程师都被精细管理，不再是 Vendor 一群人。",
)

# 第 7 页 支柱 ② 技术能力管理沉淀
make_table_page(
    page=7, total=TOTAL,
    title="支柱 ② · 技术能力管理沉淀",
    headers=["维度", "系统模块", "能力"],
    rows=[
        ["个人能力档案", "工程师 → 技能等级", "L1/L2/L3 分级，多认证挂载"],
        ["能力矩阵", "能力矩阵管理", "6 大领域热力图，一眼看团队短板"],
        ["培训记录", "培训管理", "外部认证 + 内部培训档案 + 成本字段"],
        ["IDP 个人发展", "培训 → IDP", "每人发展计划：目标 + 行动项"],
        ["知识资产库", "FDE 知识库", "7 大类资产，跨项目可复用"],
    ],
    col_widths=[3, 3.5, 5.5],
    notes="能力沉淀是长期投入的回报。系统让 团队在变强 从感觉变成数字，招标 / 评级都能拿出来证明。",
)

# 第 8 页 支柱 ③ 利润清晰可控（三口径表 + 合规高亮）
make_content(
    page=8, total=TOTAL,
    title="支柱 ③ · 利润清晰可控",
    subtitle="三种口径并存，按角色分发",
    items=[
        "A 口径 · 团队整体真实利润值（admin / lead / finance）",
        "B 口径 · 销售/客户打白工场景 + 欠款（admin / lead / finance）",
        "C 口径 · 领导层检视数据 = 服务费降本 + 打白工增效（所有人）",
    ],
    item_size=20,
    brag_box="驾驶舱视角看不到真实利润 · 三口径同时存在 · 对内透明 + 对外讲故事",
    notes="同一份原始数据出三种口径，按角色分发。对外讲降本故事，对内算真实账。这是系统的护城河。",
)

# 第 9 页 支柱 ④ 团队成本精细管理
make_content(
    page=9, total=TOTAL,
    title="支柱 ④ · 团队成本精细管理",
    subtitle="8 类支出统一抽象 · 审批流闭环",
    items=[
        "8 类支出：耗材 / 分包 / 临时人力 / 许可证 / 差旅 / 培训 / 其他 / 外包工程师支出",
        "审批流：提交 → pending → approved / rejected → paid",
        "Vendor 角色：只能提交，不能审批；只看自己 Vendor 名下数据",
        "每笔支出可追溯到项目 / 工程师 / 责任人",
        "财务一键导出报表，关账时间缩短",
    ],
    item_size=18,
    brag_box="vendor 端透明：消除「假发票 / 重复报销」的灰色空间",
    notes="成本透明化是数字化硬指标。每笔支出都有审批链可追溯。",
)

# 第 10 页 支柱 ⑤ 驾驶舱
make_table_page(
    page=10, total=TOTAL,
    title="支柱 ⑤ · 驾驶舱 · High Level 价值展示",
    headers=["Tab", "主题", "给谁看"],
    rows=[
        ["01 总览", "在管项目 + HK 地图 + 客户 Logo 墙", "老板 + 客户来访"],
        ["02 降本", "服务商 vs FDE 公司毛利率对比", "老板"],
        ["03 项目进度", "在管 / 完成 / 14 天到期 / 累计交付", "团队周会"],
        ["04 技术沉淀", "累计资产 / 跨项目复用 / 节省工时", "内部能力展示"],
        ["05 团队能力", "认证矩阵热力图 + Top 持证", "招标资质证明"],
    ],
    col_widths=[2, 6, 5],
    notes="驾驶舱是对外讲故事的窗口。客户来访可一键投屏，所有数字按合规口径展示。",
)

# 第 11 页 业务闭环
make_content(
    page=11, total=TOTAL,
    title="三、系统化落地：8 维度业务闭环",
    subtitle="立项 → 派单 → 执行 → 验收 → 复盘 → 续单 → 沉淀",
    items=[
        "立项 + 报价 + 中标 → 项目和客户管理（bid_outcome 状态机）",
        "派单 → 派单和工时管理（Assignment 接拒留痕）",
        "执行 → 工时 / 支出 / VSF（加权工时 + 审批流）",
        "验收 → 项目效率管理（对接工程师 + 互动留言）",
        "复盘 → 关键项目复盘管理（满意度 + 闭环开关）",
        "续单跟踪 → 续单跟踪（含 6 类输因，自动算胜率）",
    ],
    item_size=17,
    brag_box="工程师属于外部敏感资源，不可用 iCan 通用平台 — 需打造完美契合的专属系统",
    notes="8 环节以前散在 Excel + 邮件 + 微信群，现在串成闭环线。每个项目全部数据都在系统里。",
)

# 第 12 页 数据流动示例
make_content(
    page=12, total=TOTAL,
    title="三、一个项目的完整数字旅程",
    subtitle="所有数字自动联动 · 无需人工对账",
    items=[
        "Sales 立项 → 录入客户 + 外包估算 → bid_outcome=pending",
        "中标通知 → 改 bid_outcome=won → 驾驶舱降本数字立即跳升",
        "PM 派单 + 工程师执行 → Assignment + 加权工时 + 支出审批",
        "管理员录入中标金额 → 公司毛利率公式自动重算",
        "项目交付完成 → 写复盘 + 满意度 + 行动项闭环",
        "半年后续单 → RenewalAttempt 跟踪 → 输了记 6 类原因",
    ],
    item_size=17,
    notes="这是项目在系统里的全部生命轨迹，所有数字都自动联动。",
)

# 第 13 页 价值 · 财务
make_table_page(
    page=13, total=TOTAL,
    title="四、价值放大 · 财务维度",
    headers=["输出维度", "系统能力", "用途"],
    rows=[
        ["公司毛利率提升测算", "自动算服务商 vs FDE 两口径", "给老板讲降本故事"],
        ["单项目盈亏 (口径 B)", "团队入账 − 项目成本明细", "销售考核 / 客户分级"],
        ["团队真实利润 (口径 A)", "VSF − 全部支出", "内部预算"],
        ["驾驶舱节省金额 (口径 C)", "外包估算 − 团队入账", "对外讲故事 + 招标"],
        ["月度 / 季度趋势", "多周期对比", "战略复盘"],
    ],
    col_widths=[3.5, 4, 4.5],
    notes="财务价值不是某个具体数字，而是让 公司多赚多少 这个问题永远有据可查。",
)

# 第 14 页 价值 · 运营
make_table_page(
    page=14, total=TOTAL,
    title="四、价值放大 · 运营维度",
    headers=["运营场景", "服务商模式", "系统化后"],
    rows=[
        ["项目立项到派单", "多日邮件来回", "系统内一键完成"],
        ["月度财务关账", "Excel 多表对账", "系统报表一键导出"],
        ["工程师能力台账", "半年手工统计", "实时（季度自动快照）"],
        ["客户来访演示准备", "半天整理 PPT", "0 分钟（直接投屏驾驶舱）"],
        ["项目复盘文档", "散落各处", "100% 入库"],
    ],
    col_widths=[3.5, 4.5, 4],
    notes="运营效率提升靠机制本身的存在。以前是人在维护数据，现在是系统在维护，人只看报表。",
)

# 第 15 页 价值 · 能力
make_content(
    page=15, total=TOTAL,
    title="四、价值放大 · 能力维度",
    subtitle="系统能跟踪的团队能力指标",
    items=[
        "工程师档案：个人能力树（认证 + 技能 + IDP + 培训史）",
        "团队总数 + 在职/离职分类统计 → 驾驶舱 KPI",
        "厂商认证矩阵：6 大领域 × L1-L3，热力图 + 招标资料",
        "Top 持证工程师榜单（驾驶舱榜单）",
        "招标资质证明一键导出（L3 持证清单）",
    ],
    item_size=18,
    brag_box="客户来访可凭驾驶舱「团队能力视图」实时证明团队实力",
    notes="能力沉淀的价值在于：能用数据回答 我们团队凭什么 这个问题。招标 / 评级 / 接标杆都靠它。",
)

# 第 16 页 多角色矩阵
make_table_page(
    page=16, total=TOTAL,
    title="五、多角色：6 角色 × 11 模块",
    headers=["角色", "职责", "看什么"],
    rows=[
        ["admin", "超级管理员", "全权限，1 个内部账号"],
        ["lead", "团队负责人", "全部利润 + 战略决策"],
        ["pm", "项目经理", "立项 + 派单 + 改项目 + 审支出（不看利润）"],
        ["finance", "财务", "审支出 + 看利润 + 出报表"],
        ["engineer", "基层工程师", "只看自己派单 / 工时 / 支出"],
        ["vendor", "Vendor 联系人", "只看自己 Vendor 名下支出，不能审批"],
    ],
    col_widths=[2, 3, 7],
    notes="6 个角色 covering 内部 + 外部全部参与方。Vendor 也在系统里管，但每个人视野严格隔离。",
)

# 第 17 页 权限设计原则
make_content(
    page=17, total=TOTAL,
    title="五、权限设计 3 大原则",
    subtitle="不是配置，是架构层守门",
    items=[
        "驾驶舱合规优先：/api/cockpit/* 严禁返回 A/B 利润字段，CI 守门",
        "数据归属精确隔离：Vendor 按 vendor_id 过滤，Engineer 按 engineer_id 过滤",
        "敏感字段二次门禁：证件号明文查看需再次确认 + 写审计日志",
        "符合 PDPO（香港个人资料隐私条例）要求",
    ],
    item_size=18,
    brag_box="即使开发同事写错 API，CI 守门也会拦住 — 不可绕过的合规底线",
    notes="权限设计不是说说，是架构层就做了隔离。",
)

# 第 18 页 未来 5 方向
make_table_page(
    page=18, total=TOTAL,
    title="六、未来演进：5 大战略方向",
    headers=["#", "方向", "关键动作"],
    rows=[
        ["①", "坚持工程师外包路线", "保留人力弹性 + 不占编制；系统化让外包从黑盒变白盒"],
        ["②", "外包工程师薪资支付合规优化", "薪资透传机制 + 第三方代发 + 月度凭证审核"],
        ["③", "扩大工程师类型 + 人数", "新增 5G / 数据中心 / 售前助理 / 国际项目"],
        ["④", "系统展示更深层维度", "客户终身价值 / Vendor 性价比 / 工程师 ROI / 预测性指标"],
        ["⑤", "接更多标杆项目", "金融 / 政府 / 跨境 / IDC / 海外 — 系统化能力换大单"],
    ],
    col_widths=[0.8, 4, 8.2],
    notes="5 个方向之间是飞轮关系：更多标杆 → 更多沉淀 → 系统更深 → 更强资质 → 更多人才 → 更多标杆。",
)

# 第 19 页 结语
make_closing(page=19, total=TOTAL)


# ─── 输出 ───────────────────────────────────────────────────
from datetime import datetime
out_path = f"docs/FDE系统汇报_{datetime.now().strftime('%H%M%S')}.pptx"
prs.save(out_path)
print(f"✅ 已生成 {out_path} ({TOTAL} 页)")
